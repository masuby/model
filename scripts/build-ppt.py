#!/usr/bin/env python3
"""Generate the INFORM Tanzania subnational-risk presentation (honest narrative:
template formulas -> hidden-sheet bottleneck -> engines that mimic it -> local-data methods ->
the gap -> research & financing opportunities). Output: docs/INFORM_Tanzania_Subnational_Risk.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
NAVY = RGBColor(0x10, 0x2A, 0x43); GREEN = RGBColor(0x1E, 0x6F, 0x3D)
GOLD = RGBColor(0xF2, 0xA9, 0x00); GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); LIGHT = RGBColor(0xF1, 0xF5, 0xF9); RED = RGBColor(0xB4, 0x23, 0x18)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True; return tb.text_frame

def fill(s, l, t, w, h, color):
    sp = s.shapes.add_shape(1, l, t, w, h); sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False; return sp

def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font='Calibri'):
    p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font

def header(s, kicker, title):
    fill(s, 0, 0, SW, Inches(1.25), NAVY)
    fill(s, 0, Inches(1.25), SW, Pt(4), GOLD)
    tf = box(s, Inches(0.6), Inches(0.18), Inches(12), Inches(1.0))
    setp(tf.paragraphs[0], kicker, 12, GOLD, bold=True)
    p = tf.add_paragraph(); setp(p, title, 26, WHITE, bold=True)

def bullets(s, items, left=Inches(0.7), top=Inches(1.6), width=Inches(12), size=16, gap=6):
    tf = box(s, left, top, width, SH - top - Inches(0.4)); first = True
    for txt, lvl, color, bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        prefix = {0: '■  ', 1: '      –  ', 2: '            ·  '}[lvl]
        setp(p, prefix + txt, size - lvl*2, color, bold=bold); p.space_after = Pt(gap)

# ---- 1 TITLE ----
s = slide(); fill(s, 0, 0, SW, SH, NAVY); fill(s, 0, Inches(4.55), SW, Pt(5), GOLD)
fill(s, 0, Inches(4.6), SW, SH-Inches(4.6), GREEN)
tf = box(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(3))
setp(tf.paragraphs[0], 'INFORM TANZANIA', 14, GOLD, bold=True)
p = tf.add_paragraph(); setp(p, 'Subnational Disaster Risk', 44, WHITE, bold=True)
p = tf.add_paragraph(); setp(p, 'From a hidden-sheet Excel template to a scalable, local-data risk engine', 22, WHITE)
tf = box(s, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2))
setp(tf.paragraphs[0], "Prime Minister's Office – Disaster Management Department (PMO-DMD)", 16, WHITE, bold=True)
p = tf.add_paragraph(); setp(p, 'Reproducing the regional INFORM model exactly — and scaling it to all 195 councils', 14, RGBColor(0xD8,0xE6,0xDC))
p = tf.add_paragraph(); setp(p, 'June 2026', 13, RGBColor(0xD8,0xE6,0xDC))

# ---- 2 WHAT IS INFORM ----
s = slide(); header(s, 'BACKGROUND', 'What is INFORM?')
bullets(s, [
 ('INFORM is a global, open composite index of disaster risk — a single 0–10 score per area (0 = lowest risk, 10 = highest).', 0, NAVY, False),
 ('Tanzania runs the regional INFORM Model Template (the SADC / IGAD subnational adaptation). The official national INFORM risk is 4.1.', 0, NAVY, False),
 ('Risk is built from three dimensions — the INFORM logic:', 0, NAVY, True),
 ('Hazard & Exposure — the threats (drought, flood, earthquake, conflict…) and who/what is exposed', 1, GREY, False),
 ('Vulnerability — susceptibility of people (poverty, health, food insecurity, displacement)', 1, GREY, False),
 ('Lack of Coping Capacity — weakness of institutions & infrastructure to respond (WASH, health, DRR, governance)', 1, GREY, False),
 ('Each dimension is built from indicators → categories → dimension → the final risk index.', 0, GREEN, True),
], top=Inches(1.7), size=17)

# ---- 3 TEMPLATE STRUCTURE ----
s = slide(); header(s, 'THE MODEL TEMPLATE', 'Structure: 78 indicators → 6 categories → 3 dimensions → 1 risk')
cols = [('HAZARD & EXPOSURE', ['Natural (drought, flood, earthquake,', 'landslide, cyclone, coastal, wildfire…)', 'Human (conflict, violence)'], GREEN),
        ('VULNERABILITY', ['Socio-economic (poverty, GDP,', 'aid dependency, habitat)', 'Vulnerable groups (health, nutrition,', 'food security, displaced people)'], GOLD),
        ('LACK OF COPING CAPACITY', ['Institutional (DRR, governance)', 'Infrastructure (WASH, health access,', 'communication, education)'], RED)]
x = Inches(0.6)
for title, lines, col in cols:
    w = Inches(3.9); fill(s, x, Inches(1.8), w, Inches(0.7), col)
    tf = box(s, x, Inches(1.86), w, Inches(0.6)); setp(tf.paragraphs[0], title, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
    bx = fill(s, x, Inches(2.5), w, Inches(3.4), LIGHT); bx.line.color.rgb = col; bx.line.width = Pt(1)
    tf = box(s, x+Inches(0.15), Inches(2.65), w-Inches(0.3), Inches(3.2)); fr = True
    for ln in lines:
        p = tf.paragraphs[0] if fr else tf.add_paragraph(); fr = False; setp(p, ln, 13, NAVY); p.space_after = Pt(6)
    x += Inches(4.0)
tf = box(s, Inches(0.6), Inches(6.2), Inches(12), Inches(1))
setp(tf.paragraphs[0], '78 raw indicators (53 currently used for Tanzania) → aggregated to 6 categories → 3 dimensions → one INFORM risk score, for every administrative unit.', 15, GREEN, bold=True)

# ---- 4 INDICATOR FORMULA CHAIN ----
s = slide(); header(s, 'HOW THE TEMPLATE WORKS', 'From an actual value to a 0–10 score — the standardisation chain')
steps = ['Actual value (natural unit)', 'Denominator (÷ population / area / GDP, or none)',
         'Outlier cap — Tukey IQR fence  MAX(MIN(x, Q3+1.5·IQR), Q1−1.5·IQR)',
         'Transform — None | Logarithm LN(0.001+x) | Exponential',
         'Min–max scaling — 10 × (x − Min) / (Max − Min)  vs a fixed reference',
         'Direction — if protective: 10 − score (higher = better)',
         'Clamp [0,10] → ROUND to 1 dp →  the 0–10 indicator']
y = Inches(1.7)
for i, st in enumerate(steps):
    col = GREEN if i == len(steps)-1 else NAVY
    bx = fill(s, Inches(0.7), y, Inches(11.9), Inches(0.6), LIGHT if i < len(steps)-1 else RGBColor(0xE6,0xF2,0xE9))
    bx.line.color.rgb = col; bx.line.width = Pt(1)
    tf = box(s, Inches(0.9), y+Inches(0.06), Inches(11.5), Inches(0.5))
    setp(tf.paragraphs[0], f'{i+1}.  {st}', 14, col, bold=(i==len(steps)-1)); tf.paragraphs[0].space_after = Pt(0)
    y += Inches(0.72)
tf = box(s, Inches(0.7), y+Inches(0.02), Inches(12), Inches(0.5))
setp(tf.paragraphs[0], 'Excel:  =IF(val="","",IFERROR(MAX(0,MIN(10,ROUND(IF(sign="Decrease",10-base,base),1))),"No data"))', 12, GREY, bold=True, font='Consolas')

# ---- 5 FEW INDICATORS, MANY UNITS ----
s = slide(); header(s, 'CAPTURING DIVERSITY', 'A few indicators — many different units (why standardisation is needed)')
rows = [('Indicator', 'Unit', 'Dimension', True),
        ('Food security – insufficient food', '% of population (Phase 2+)', 'Vulnerability', False),
        ('Historic drought frequency', 'years', 'Hazard & Exposure', False),
        ('Internally displaced people', 'count (number of people)', 'Vulnerability', False),
        ('Gross National Income per capita', 'US$ (thousands, log)', 'Coping Capacity', False),
        ('Health expenditure per capita', 'US$ per capita', 'Coping Capacity', False),
        ('Physicians density', 'per 10,000 population', 'Coping Capacity', False),
        ('Measles incidence', 'per 1,000,000 population', 'Vulnerability', False),
        ('Soil erosion', 'Mg / ha / yr', 'Hazard & Exposure', False),
        ('Coastal erosion', 'metres / year', 'Hazard & Exposure', False)]
tw = Inches(12); th = Inches(4.9); tbl = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.7), tw, th).table
tbl.columns[0].width = Inches(5.2); tbl.columns[1].width = Inches(4.0); tbl.columns[2].width = Inches(2.8)
for ri, (a, b, c, hd) in enumerate(rows):
    for ci, val in enumerate((a, b, c)):
        cell = tbl.cell(ri, ci); cell.text = val
        para = cell.text_frame.paragraphs[0]; para.runs[0].font.size = Pt(15 if hd else 13)
        para.runs[0].font.bold = hd; para.runs[0].font.color.rgb = WHITE if hd else NAVY
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if hd else (LIGHT if ri % 2 else WHITE)
tf = box(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.5))
setp(tf.paragraphs[0], '%, years, counts, US$, per-1,000, per-10,000, ratios, Mg/ha/yr, m/yr — all reconciled onto one 0–10 scale.', 14, GREEN, bold=True)

# ---- 6 CATEGORISATION -> RISK ----
s = slide(); header(s, 'CATEGORISATION & AGGREGATION', 'From standardised indicators to the INFORM risk')
bullets(s, [
 ('Component  =  AVERAGE of its indicators where Use = "Yes"   (Excel AVERAGEIFS)', 0, NAVY, True),
 ('Category   =  AVERAGE of its components   (arithmetic mean)', 0, NAVY, True),
 ('Dimension  =  scaled GEOMEAN of its categories', 0, NAVY, True),
 ('ROUND( (10 − GEOMEAN((10−c)/10×9+1)) /9 ×10 , 1 )', 1, GREY, False),
 ('INFORM Risk  =  ∛ ( H × V × LCC )   — cube-root of the three dimensions', 0, GREEN, True),
 ('ROUND( H^(1/3) × V^(1/3) × LCC^(1/3) , 1 )', 1, GREY, False),
 ('The geometric mean means a country cannot "average away" one very weak dimension — high hazard with low coping still yields high risk.', 0, NAVY, False),
], top=Inches(1.8), size=18, gap=10)

# ---- 7 THE BOTTLENECK ----
s = slide(); header(s, 'THE HONEST PROBLEM', 'The bottleneck: it all lives in deeply hidden Excel sheets')
fill(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.7), RED)
tf = box(s, Inches(0.9), Inches(1.74), Inches(11.6), Inches(0.6))
setp(tf.paragraphs[0], 'Every reference range, outlier fence, transform and aggregation is buried in hidden worksheets.', 16, WHITE, bold=True)
bullets(s, [
 ('Opaque — the math is invisible to the people who use and must trust the numbers.', 0, NAVY, False),
 ('Fixed unit set — built for ~170 districts; it cannot scale to 195 councils, wards or villages.', 0, NAVY, False),
 ('Manual & single-file — one workbook, no version control, easy to break, hard to audit.', 0, NAVY, False),
 ('Closed to local data — no clean way to add a local indicator or feed council-level values.', 0, NAVY, False),
 ('No automation — no live update, no link to early-warning or anticipatory action.', 0, NAVY, False),
 ('Result: the model stops at the national / regional level. Local risk — where decisions are made — is out of reach.', 0, RED, True),
], top=Inches(2.6), size=16, gap=11)

# ---- 8 ENGINES MIMIC EXCEL ----
s = slide(); header(s, 'THE ENABLER', 'We rebuilt the hidden sheets as a transparent engine — proven identical')
bullets(s, [
 ('The hidden-sheet formulas were reverse-engineered and reproduced exactly in code (a standardisation + aggregation engine).', 0, NAVY, False),
 ('Proven byte-for-byte against the official workbook:', 0, GREEN, True),
 ('8,664 / 8,664 standardised indicator cells reproduce the template exactly', 1, GREEN, False),
 ('170 / 170 district risk scores reproduce the template exactly', 1, GREEN, False),
 ('Now it scales — the same engine runs unchanged on 195 councils and 31 regions (194/195 match district INFORM; the rest are new post-2024 councils).', 0, NAVY, False),
 ('Spec-driven & open — add, delete or amend an indicator without touching the math; fully testable and version-controlled.', 0, NAVY, False),
 ('The regional INFORM core is never altered — it is reproduced exactly; everything new is additive (v2).', 0, GOLD, True),
], top=Inches(1.7), size=16, gap=10)

# ---- 9 LOCAL DATA METHODS ----
s = slide(); header(s, 'CAPTURING REAL LOCAL DATA', 'Methods — a federated, bottom-up data system')
bullets(s, [
 ('Federated collection tools — one shared Excel sent to every sector; they enter ACTUAL values in natural units and a locked formula shows the 0–10 live.', 0, NAVY, False),
 ('One accountable Tanzanian stakeholder per indicator:', 0, NAVY, True),
 ('TMA (climate), GST (geology), NEMC (environment), MoH (health), NBS (statistics), MoA/MUCHALI (food), MoW (water), PMO-DMD (DRR)…', 1, GREY, False),
 ('Bottom-up aggregation — village → district → region → country; finer, well-filled data refines the higher level.', 0, NAVY, False),
 ('Authentic sources used: NBS 2022 Census, TMA CHIRPS v3 / ERA5, GST seismic catalogue, MoH TDHS-MIS 2022, IPC/MUCHALI, TASAF, and satellite Earth Observation.', 0, NAVY, False),
 ('Earth Observation computed live (Earth Engine): MODIS NDVI, SMAP soil moisture, CHIRPS aridity — each validated against known geography before use.', 0, GREEN, True),
], top=Inches(1.7), size=16, gap=10)

# ---- 9b HONEST STATUS (for model developers) ----
s = slide(); header(s, 'HONEST STATUS', 'Where this actually stands — no overselling')
quad = [('SOLID / EXACT', GREEN, [
          'Engine reproduces the template byte-for-byte: 8,664/8,664 cells, 170/170 risks',
          'Same engine runs at 195 councils / 31 regions (194/195 = district INFORM)',
          '2022 census population per council (sum = 61,741,120, the exact national total)']),
        ('PROTOTYPE / PROPOSED', GOLD, [
          'Advanced "exploded" multi-source indicators + weights = literature-proposed, NOT enforced',
          'EO drought layers (aridity P/PET, NDVI, soil moisture) computed & validated — research-grade',
          'Seasonality from CHIRPS 8-zone analysis — context layer, not yet in the index']),
        ('TRIED — DID NOT WORK', RED, [
          'Sentinel-1 flood (VV<-17dB): FAILED validation — flagged soda-lakes/salt-flats as water → PULLED',
          'Lightning (LIS/OTD): not available on Earth Engine — not done',
          'SMAP soil moisture product is deprecated (ends 2022) — not current']),
        ('STILL OPEN', NAVY, [
          'Many vulnerability indicators exist only at region/national level → council = labelled proxy',
          'No sustainable update pipeline; no live link to EOCC / early warning',
          'TASAF council poverty & IPC food: data exists but locked in PDFs/registry'])]
xy = [(Inches(0.55), Inches(1.65)), (Inches(6.85), Inches(1.65)), (Inches(0.55), Inches(4.55)), (Inches(6.85), Inches(4.55))]
for (title, col, items), (x, y) in zip(quad, xy):
    w = Inches(5.9); h = Inches(2.75)
    fill(s, x, y, w, Pt(28), col)
    tf = box(s, x+Inches(0.1), y+Pt(2), w-Inches(0.2), Pt(26)); setp(tf.paragraphs[0], title, 13, WHITE, bold=True)
    bx = fill(s, x, y+Pt(28), w, h-Pt(28), LIGHT); bx.line.color.rgb = col; bx.line.width = Pt(1)
    tf = box(s, x+Inches(0.18), y+Pt(34), w-Inches(0.36), h-Pt(40)); fr = True
    for it in items:
        p = tf.paragraphs[0] if fr else tf.add_paragraph(); fr = False
        setp(p, '·  ' + it, 11, NAVY); p.space_after = Pt(5)

# ---- 10 THE GAP ----
s = slide(); header(s, 'THE GAP', 'What is still missing — the challenge ahead')
bullets(s, [
 ('Subnational risk is not yet operational — decisions happen at council/ward level, but official INFORM stops at national/regional.', 0, RED, True),
 ('Local data is scattered and uneven — many key indicators exist only at region or national level (poverty, health surveys, GDP).', 0, NAVY, False),
 ('No sustainable pipeline — capturing, validating and updating local data is still manual and ad-hoc.', 0, NAVY, False),
 ('Earth-Observation indicators (vegetation, soil moisture, flood, landslide) need proper computation, downscaling and validation.', 0, NAVY, False),
 ('No live link to action — risk is not yet wired to the EOCC, early-warning, or LGA planning & budgeting.', 0, NAVY, False),
 ('Capacity & financing gaps at the local-government level — tools, training and incentives to keep the data flowing.', 0, NAVY, False),
], top=Inches(1.7), size=16, gap=12)

# ---- 11 OPPORTUNITIES ----
s = slide(); header(s, 'OPPORTUNITIES', 'Research & financial support — for member-state adoption')
two = [('RESEARCH', GREEN, [
        'EO computation & statistical downscaling to council/ward',
        'Validation against ground events & independent products',
        'Multi-source "exploded" indicators (one hazard, many institutions)',
        'Climate seasonality (8 rainfall zones) into risk & early warning',
        'Methods publishable as regional public goods']),
       ('FINANCING & ADOPTION', GOLD, [
        'Fund local-data integration & the federated tools',
        'LGA capacity building — train sectors to capture & submit',
        'Scale to SADC / IGAD member states (transparent, open engine)',
        'Wire risk to PMO-DMD EOCC, anticipatory action & PO-RALG planning',
        'Tanzania as the subnational-INFORM pilot for the region'])]
x = Inches(0.6)
for title, col, items in two:
    w = Inches(6.0); fill(s, x, Inches(1.75), w, Inches(0.7), col)
    tf = box(s, x, Inches(1.82), w, Inches(0.6)); setp(tf.paragraphs[0], title, 16, WHITE, bold=True, align=PP_ALIGN.CENTER)
    bx = fill(s, x, Inches(2.45), w, Inches(4.2), LIGHT); bx.line.color.rgb = col; bx.line.width = Pt(1.5)
    tf = box(s, x+Inches(0.25), Inches(2.65), w-Inches(0.5), Inches(3.9)); fr = True
    for it in items:
        p = tf.paragraphs[0] if fr else tf.add_paragraph(); fr = False
        setp(p, '■  ' + it, 15, NAVY); p.space_after = Pt(12)
    x += Inches(6.4)

# ---- 12 CLOSING ----
s = slide(); fill(s, 0, 0, SW, SH, NAVY); fill(s, 0, Inches(3.0), SW, Pt(5), GOLD)
tf = box(s, Inches(0.9), Inches(1.4), Inches(11.6), Inches(2))
setp(tf.paragraphs[0], 'Summary', 30, WHITE, bold=True)
tf = box(s, Inches(0.9), Inches(3.3), Inches(11.6), Inches(3))
for i, t in enumerate([
   'The template math is fully understood and reproduced exactly (8,664/8,664; 170/170).',
   'It runs at council level (195) because it is code, not a hidden-sheet workbook.',
   'Local-data capture and EO are working but research-grade — and some attempts failed (stated plainly).',
   'Ask: research partnership + financing to integrate local data and support member-state adoption.']):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    setp(p, '▸  ' + t, 18, RGBColor(0xD8,0xE6,0xDC)); p.space_after = Pt(12)
tf = box(s, Inches(0.9), Inches(6.6), Inches(11.6), Inches(0.6))
setp(tf.paragraphs[0], "PMO-DMD Tanzania   ·   INFORM Subnational Risk   ·   June 2026", 13, GOLD, bold=True)

out = os.path.join(ROOT, 'docs/INFORM_Tanzania_Subnational_Risk.pptx')
prs.save(out)
print('wrote', out, '|', len(prs.slides._sldIdLst), 'slides')
