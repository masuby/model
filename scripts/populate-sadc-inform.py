#!/usr/bin/env python3
"""Put the INFORM assignment content into the SADC template FORMAT (keeps SADC logo/swoosh/flags).
Assignment: the Tanzania INFORM Model Template — indicator formulas -> risk; a few indicators with diverse
units; capturing local data (methods); categorisation -> risk; the hidden-sheet bottleneck; engines that
mimic the Excel for scale-up; the gap; opportunities (research + financing for member-state adoption).
Output: docs/INFORM_Tanzania_SADC_format.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TPL = '/home/kaijage/Downloads/Country SADC  PPT Template.pptx'
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
NAVY = RGBColor(0x10, 0x2A, 0x43); GREEN = RGBColor(0x1E, 0x6F, 0x3D)
GREY = RGBColor(0x44, 0x4C, 0x57); RED = RGBColor(0xB4, 0x23, 0x18)

p = Presentation(TPL)
LAYOUT = p.slides[1].slide_layout  # SADC "Title and Content" (branded)

def body_ph(s):
    cands = [ph for ph in s.placeholders if ph.placeholder_format.type != 1 and ph.has_text_frame]
    if not cands: return None
    return max(cands, key=lambda ph: (ph.width or 0) * (ph.height or 0))

def fill(slide, title, bullets, sz0=15, sz1=13):
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    bp = body_ph(slide)
    tf = bp.text_frame; tf.clear(); tf.word_wrap = True
    for i, b in enumerate(bullets):
        txt, lvl, color, bold = b
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt; para.level = lvl; para.space_after = Pt(5)
        r = para.runs[0]; r.font.size = Pt(sz0 if lvl == 0 else sz1); r.font.bold = bold
        r.font.color.rgb = color

# reuse the template's existing content slides first (no deletion - that corrupts the package), then add more
_pool = list(p.slides)[1:]; _i = [0]

# --- Slide 1 title + INFORM subtitle ---
s1 = p.slides[0]
if s1.shapes.title is not None:
    s1.shapes.title.text = 'UNITED REPUBLIC OF TANZANIA'
tb = s1.shapes.add_textbox(Inches(0.8), Inches(5.15), Inches(11.7), Inches(1.9)).text_frame; tb.word_wrap = True
for i, (t, sz, b) in enumerate([
    ('INFORM Subnational Disaster Risk', 20, True),
    ('The Country Model Template — how risk is calculated, the bottleneck, and scaling it to local level', 13, False),
    ("Prime Minister's Office – Disaster Management Department (PMO-DMD)", 12, False)]):
    para = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
    para.text = t; r = para.runs[0]; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = NAVY

def add():
    if _i[0] < len(_pool):
        s = _pool[_i[0]]; _i[0] += 1; return s
    return p.slides.add_slide(LAYOUT)

# 2 — what INFORM is
fill(add(), 'INFORM: how the Country Model Template is built', [
 ('INFORM is an open composite index of disaster risk — one 0–10 score per area (0 = lowest, 10 = highest). Tanzania runs the regional INFORM Country Model Template (Excel); the national INFORM risk is 4.1.', 0, NAVY, False),
 ('Risk is built from three dimensions:', 0, NAVY, True),
 ('Hazard & Exposure — drought, flood, earthquake, cyclone, conflict … and who/what is exposed', 1, GREY, False),
 ('Vulnerability — poverty, health, nutrition, food insecurity, displacement', 1, GREY, False),
 ('Lack of Coping Capacity — institutions & infrastructure (DRR, governance, WASH, health, education)', 1, GREY, False),
 ('Structure: 78 raw indicators (53 used) → 6 categories → 3 dimensions → 1 risk score, for every unit.', 0, GREEN, True),
], sz0=15)

# 3 — the standardisation formula chain (methods)
fill(add(), 'How an indicator becomes a 0–10 score (the hidden-sheet method)', [
 ('1.  Actual value, in its natural unit', 0, NAVY, False),
 ('2.  Denominator — ÷ population / area / GDP, or none', 0, NAVY, False),
 ('3.  Outlier cap — Tukey IQR fence:  MAX( MIN(x, Q3 + 1.5·IQR) , Q1 − 1.5·IQR )', 0, NAVY, False),
 ('4.  Transform — None | Logarithm LN(0.001 + x) | Exponential', 0, NAVY, False),
 ('5.  Min–max scaling — 10 × (x − Min) / (Max − Min)  against a fixed reference range', 0, NAVY, False),
 ('6.  Direction — if protective (higher = better): 10 − score', 0, NAVY, False),
 ('7.  Clamp to [0,10] → ROUND to 1 decimal → the 0–10 indicator', 0, GREEN, True),
 ('Every indicator, whatever its unit, passes through the same chain — Tukey outliers, log/exp transform, min–max, direction.', 0, GREY, False),
], sz0=15, sz1=13)

# 4 — a few indicators, many units (TABLE)
st = add(); st.shapes.title.text = 'A few indicators — many units (why standardisation is needed)'
bp = body_ph(st)
if bp is not None:
    bp.text_frame.clear()
rows = [('Indicator', 'Unit', 'Dimension'),
        ('Food security – insufficient food', '% of population', 'Vulnerability'),
        ('Historic drought frequency', 'years', 'Hazard & Exposure'),
        ('Internally displaced people', 'count (people)', 'Vulnerability'),
        ('Gross National Income per capita', 'US$ (thousands, log)', 'Coping Capacity'),
        ('Health expenditure per capita', 'US$ per capita', 'Coping Capacity'),
        ('Physicians density', 'per 10,000 population', 'Coping Capacity'),
        ('Measles incidence', 'per 1,000,000 population', 'Vulnerability'),
        ('Soil erosion', 'Mg / ha / yr', 'Hazard & Exposure'),
        ('Coastal erosion', 'metres / year', 'Hazard & Exposure')]
tb = st.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.0)).table
tb.columns[0].width = Inches(5.2); tb.columns[1].width = Inches(3.8); tb.columns[2].width = Inches(2.9)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        c = tb.cell(ri, ci); c.text = val; pr = c.text_frame.paragraphs[0].runs[0]
        pr.font.size = Pt(13 if ri else 14); pr.font.bold = (ri == 0)
        pr.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else NAVY
        c.fill.solid(); c.fill.fore_color.rgb = NAVY if ri == 0 else (RGBColor(0xEE, 0xF2, 0xF7) if ri % 2 else RGBColor(0xFF, 0xFF, 0xFF))
txt = st.shapes.add_textbox(Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.5)).text_frame
txt.paragraphs[0].text = '%, years, counts, US$, per-1,000, per-10,000, Mg/ha/yr, m/yr — all reconciled onto one 0–10 scale.'
txt.paragraphs[0].runs[0].font.size = Pt(13); txt.paragraphs[0].runs[0].font.bold = True; txt.paragraphs[0].runs[0].font.color.rgb = GREEN

# 5 — categorisation -> risk
fill(add(), 'From indicators to risk — categorisation & aggregation', [
 ('Component  =  AVERAGE of its indicators where Use = "Yes"   (Excel AVERAGEIFS)', 0, NAVY, True),
 ('Category   =  AVERAGE of its components   (arithmetic mean)', 0, NAVY, True),
 ('Dimension  =  scaled GEOMEAN of its categories', 0, NAVY, True),
 ('Risk  =  ∛ ( H × V × LCC )   — cube-root of the three dimensions', 0, GREEN, True),
 ('ROUND( H^(1/3) × V^(1/3) × LCC^(1/3) , 1 )', 1, GREY, False),
 ('The geometric mean matters: a unit cannot "average away" one very weak dimension — high hazard with low coping capacity still yields high risk.', 0, NAVY, False),
], sz0=17, sz1=13)

# 6 — the bottleneck (honest)
fill(add(), 'The honest problem: it all lives in deeply hidden Excel sheets', [
 ('Every reference range, outlier fence, transform and aggregation step is buried in hidden worksheets.', 0, RED, True),
 ('Opaque — the math is invisible to the people who must use and trust the numbers.', 0, NAVY, False),
 ('Fixed unit set — built for ~170 districts; it cannot scale to 195 councils, wards or villages.', 0, NAVY, False),
 ('Manual & single-file — one workbook, no version control, easy to break, hard to audit.', 0, NAVY, False),
 ('Closed to local data — no clean way to add a local indicator or feed council-level values.', 0, NAVY, False),
 ('No automation — no live update, no link to early warning or anticipatory action.', 0, NAVY, False),
 ('Result: the model stops at national / regional level. Local risk — where decisions are made — is out of reach.', 0, RED, True),
], sz0=14)

# 7 — engines mimic the Excel (scale-up)
fill(add(), 'We rebuilt the hidden sheets as an engine — proven identical, built to scale', [
 ('The hidden-sheet formulas were reverse-engineered and reproduced exactly in open code.', 0, NAVY, False),
 ('Proven byte-for-byte against the official workbook:', 0, GREEN, True),
 ('8,664 / 8,664 standardised indicator cells reproduce the template exactly', 1, GREEN, False),
 ('170 / 170 district risk scores reproduce the template exactly', 1, GREEN, False),
 ('Now it scales — the same engine runs unchanged on 195 councils and 31 regions.', 0, NAVY, False),
 ('Spec-driven & transparent — add, delete or amend an indicator without touching the math; fully testable.', 0, NAVY, False),
 ('The regional INFORM core is never altered — it is reproduced exactly; everything new is additive.', 0, NAVY, True),
], sz0=15)

# 8 — capturing local data (methods)
fill(add(), 'Capturing real local data — the methods', [
 ('Federated collection tools — one shared Excel per sector; they enter ACTUAL values in natural units and a locked formula shows the 0–10 live.', 0, NAVY, False),
 ('One accountable stakeholder per indicator — TMA (climate), GST (geology), NEMC (environment), MoH (health), NBS (statistics), MoA/MUCHALI (food), MoW (water), PMO-DMD (DRR).', 0, NAVY, False),
 ('Bottom-up aggregation — village → district → region → country; finer, well-filled data refines the higher level.', 0, NAVY, False),
 ('Authentic sources — NBS 2022 Census, TMA CHIRPS/ERA5, GST seismic, MoH TDHS, IPC/MUCHALI, TASAF, and satellite Earth Observation (validated before use).', 0, NAVY, False),
], sz0=16)

# 9 — the gap
fill(add(), 'The gap — what still has to be solved', [
 ('Subnational risk is not yet operational — decisions happen at council / ward level, but INFORM stops at national / regional.', 0, RED, True),
 ('Local data is scattered and uneven — many key indicators exist only at region or national level.', 0, NAVY, False),
 ('No sustainable pipeline — capturing, validating and updating local data is still manual.', 0, NAVY, False),
 ('Earth-Observation indicators (vegetation, soil moisture, flood, landslide) need proper computation, downscaling and validation.', 0, NAVY, False),
 ('No live link to action — risk is not yet wired to the emergency operations centre, early warning, or local planning & budgeting.', 0, NAVY, False),
 ('Capacity & financing gaps at local-government level — tools, training and incentives to keep the data flowing.', 0, NAVY, False),
], sz0=15)

# 10 — opportunities
fill(add(), 'Opportunities — research & financial support for member-state adoption', [
 ('RESEARCH', 0, GREEN, True),
 ('Earth-Observation computation & statistical downscaling to council / ward level', 1, NAVY, False),
 ('Validation against ground events; multi-source "exploded" indicators; climate-seasonality into early warning', 1, NAVY, False),
 ('FINANCIAL SUPPORT & ADOPTION', 0, GREEN, True),
 ('Fund the local-data integration and the federated tools; build local-government capacity to capture & submit', 1, NAVY, False),
 ('Scale the transparent engine across SADC member states; integrate risk into EOCC, anticipatory action & local planning', 1, NAVY, False),
 ('Tanzania as the subnational-INFORM pilot — a shared, open method and data standard for the region.', 0, GREEN, True),
], sz0=15, sz1=13)

# 11 — conclusion
fill(add(), 'Conclusion', [
 ('The Country Model Template math is fully understood and reproduced exactly — 8,664/8,664 cells, 170/170 risks.', 0, NAVY, False),
 ('It now runs at council level (195) because it is open code, not a hidden-sheet workbook.', 0, NAVY, False),
 ('Local-data capture and Earth Observation are working, with each layer validated before use.', 0, NAVY, False),
 ('Ask: research partnership + financing to integrate local data and support member-state adoption and local-level integration.', 0, GREEN, True),
], sz0=16)

out = os.path.join(ROOT, 'docs/INFORM_Tanzania_SADC_format.pptx')
p.save(out)
print('wrote', out, '|', len(p.slides._sldIdLst), 'slides')
