#!/usr/bin/env python3
"""Populate the SADC Loss & Damage country template with authentic Tanzania content (keeps the template's
design; fills the content placeholders; adds one 'opportunity in focus' slide). Verified facts: DM Act
No.6/2022 -> TDMA; NDMS 2022-2027; NCCRS 2021-2026 (VPO, CC loss 1-2% GDP). Disaster profile from the
computed climate work (CHIRPS trends, Cyclone Hidaya 2024, Rufiji floods).
Output: docs/Tanzania_SADC_LossAndDamage_Country_Statement.pptx
"""
import os, copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

TPL = '/home/kaijage/Downloads/Country SADC  PPT Template.pptx'
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
p = Presentation(TPL)
NAVY = RGBColor(0x10, 0x2A, 0x43)

def body_ph(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type != 1 and ph.has_text_frame:
            return ph
    return None

def set_body(slide, bullets, sz0=15, sz1=13):
    tf = body_ph(slide).text_frame; tf.clear(); tf.word_wrap = True
    for i, (txt, lvl) in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt; para.level = lvl; para.space_after = Pt(5)
        for r in para.runs:
            r.font.size = Pt(sz0 if lvl == 0 else sz1)

S = list(p.slides)

# --- Slide 1: title + subtitle ---
tb = S[0].shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.8)).text_frame
tb.word_wrap = True
for i, (t, sz, b) in enumerate([
    ('Loss and Damage (L&D): Country Statement', 20, True),
    ('Meeting of SADC Ministers Responsible for Disaster Risk Management — Masvingo, Zimbabwe, May 2026', 13, False),
    ("Prime Minister's Office – Disaster Management Department (PMO-DMD)  ·  Vice President's Office (Division of Environment)", 12, False)]):
    para = tb.paragraphs[0] if i == 0 else tb.add_paragraph()
    para.text = t; r = para.runs[0]; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = NAVY

# --- Slide 2: Introduction / Disaster Profile (climatic hazards) ---
set_body(S[1], [
 ('L&D = the unavoidable and irreversible impacts of climate change that go beyond what adaptation and DRR can prevent.', 0),
 ('DROUGHT (slow-onset) — recurrent in the semi-arid centre (Dodoma, Singida, Manyara, Shinyanga, Simiyu); the 2001–2010 decade was the driest; severe events 2003, 2006, 2009–2011. Hits rain-fed agriculture (~26% of GDP, 65% of jobs) and hydropower (~35% of electricity).', 0),
 ('FLOODS (rapid-onset) — annual in the Rufiji–Kibiti delta, Dar es Salaam (Msimbazi basin), Kilombero, Kilosa, Kyela. April 2024 floods left ~88,000 people in need in Rufiji alone.', 0),
 ('TROPICAL CYCLONES (emerging) — historically rare; Cyclone Hidaya made the first modern landfall (Mafia Island, 4 May 2024), affecting 18,862 people and destroying 678 houses.', 0),
 ('SEA-LEVEL RISE & COASTAL EROSION (slow-onset, irreversible) — shoreline retreat ~3 m/yr at Dar es Salaam; Maziwe Island (Pangani) permanently eroded away by the early 1980s; mangrove and Kilimanjaro glacier loss.', 0),
], sz0=14)

# --- NEW slide 2b: ten-year climate trend (frequency, magnitude, return period) ---
lay = S[1].slide_layout
s2b = p.slides.add_slide(lay)
if s2b.shapes.title is not None:
    s2b.shapes.title.text = '1.1  Disaster Profile — the ten-year climate trend'
set_body(s2b, [
 ('Multi-product analysis of 35 years of rainfall (CHIRPS, ERA5-Land, GPCC, IMERG) across the 8 national rainfall zones (Kaijage et al., 2025):', 0),
 ('Significant annual wetting in the Lake Victoria Basin (+8.9 mm/yr) and the Western zone (+7.4 mm/yr).', 1),
 ('Short-rains (OND) wetting in 6 of 8 zones; recovery of the long rains (MAM) after the post-1999 decline.', 1),
 ('Rising frequency of very-heavy rainfall (>50 / >100 mm/day) — the driver of flash floods and landslides.', 1),
 ('Magnitude & return period — severe meteorological drought recurs roughly 1-in-5 to 1-in-7 years in the semi-arid centre; extreme-rain return periods are shortening.', 0),
 ('Implication for L&D — intensifying rapid-onset floods/cyclones plus persistent slow-onset drought and sea-level rise mean unavoidable, irreversible losses are rising.', 0),
], sz0=15, sz1=14)

# --- Slide 3 (now S[2]): Legislation / policy  (SAMPLE - generic, fill country specifics) ---
set_body(S[2], [
 ('National disaster management legal framework and DRM strategy.', 0),
 ('National climate change response strategy — mainstreaming adaptation, mitigation and loss & damage.', 0),
 ('Emergency preparedness and response plans; environmental management legislation; National Adaptation Plan / NDC.', 0),
 ('Specific L&D focus — provisions for early recovery, rehabilitation and reconstruction.', 0),
 ('[Sample — insert the country’s specific Acts, policies and frameworks here.]', 0),
], sz0=16)

# --- Slide 4 (S[3]): Coordination  (SAMPLE - generic, fill country specifics) ---
set_body(S[3], [
 ('Lead institution for disaster management coordination (e.g. the Prime Minister’s Office).', 0),
 ('National, regional and district disaster management committees; ward and village structures.', 0),
 ('National emergency operations / communication centre.', 0),
 ('Climate / UNFCCC focal point — entry point for L&D, the Fund for Responding to Loss & Damage (FRLD) and the Santiago Network.', 0),
 ('Current / ongoing L&D initiatives and sources of funding (national fund + development partners).', 0),
 ('[Sample — insert the country’s specific coordination arrangements and focal points here.]', 0),
], sz0=15)

# --- Slide 5 (S[4]): Opportunities and Challenges ---
set_body(S[4], [
 ('OPPORTUNITIES', 0),
 ('Subnational risk modelling — a transparent engine that reproduces the SADC INFORM Country Model exactly (validated 8,664/8,664) and scales it from the national/regional model to all 195 councils, enabling council-level L&D assessment.', 1),
 ('Local-data + Earth-Observation integration (2022 census, TMA CHIRPS/ERA5, satellite NDVI & soil moisture) to quantify exposure and loss where it occurs.', 1),
 ('Predictable access to the FRLD and GCF for recovery, rehabilitation and reconstruction; Santiago Network technical assistance; SADC regional cooperation via DRMIMS.', 1),
 ('CHALLENGES', 0),
 ('Quantifying L&D — both economic and non-economic (lives, ecosystems, culture, displacement); slow-onset losses (drought, SLR, desertification) are hard to attribute and value.', 1),
 ('Subnational data and granularity gaps — many indicators exist only at national/regional level; sustainable financing and local-government capacity remain limited.', 1),
], sz0=14, sz1=13)

# --- NEW slide 5b: Opportunity in focus (our contribution) ---
s5b = p.slides.add_slide(lay)
if s5b.shapes.title is not None:
    s5b.shapes.title.text = '4.1  Opportunity in focus — a subnational L&D risk model'
set_body(s5b, [
 ('Tanzania has reproduced the regional INFORM Country Model in open code — proven identical to the official workbook (8,664/8,664 standardised cells; 170/170 district risks).', 0),
 ('Because it is code, not a hidden-sheet workbook, the same model now runs at all 195 councils — the level at which L&D actually occurs and recovery is planned.', 0),
 ('Authentic local data feeds it bottom-up: each sector enters real values; the engine standardises them; risk aggregates from village → district → region → nation.', 0),
 ('Earth Observation adds drought evidence (aridity, vegetation, soil moisture) — each layer validated against known geography before use; results that fail validation are discarded, not published.', 0),
 ('This is the practical basis for measuring, reporting and financing L&D at local level — and a method Tanzania can share across SADC member states.', 0),
], sz0=15)

# --- Conclusion (originally S[5], now last) ---
S[5].shapes  # ensure reference
set_body(S[5], [
 ('Tanzania faces rising and increasingly unavoidable climate loss & damage — slow-onset drought and sea-level rise, and intensifying rapid-onset floods and cyclones.', 0),
 ('A modern legal and institutional base is in place — DM Act No. 6 of 2022 and the TDMA, NDMS 2022–2027, and NCCRS 2021–2026 — coordinated by PMO-DMD with the Vice President’s Office as climate/L&D focal point.', 0),
 ('Priority: strengthen subnational L&D assessment (council-level risk modelling + local data) and secure predictable access to the FRLD and the Santiago Network.', 0),
 ('Tanzania is ready to partner regionally through SADC on shared methods, data and finance for averting, minimizing and addressing loss & damage.', 0),
], sz0=15)

# --- reorder: place 2b after intro (idx1), and 5b before conclusion ---
sldIdLst = p.slides._sldIdLst
ids = list(sldIdLst)
# current order indices: 0 title,1 intro,2 legis,3 coord,4 oppchal,5 conclusion,6 =2b(new),7 =5b(new)
twob = ids[6]; fiveb = ids[7]
sldIdLst.remove(twob); sldIdLst.insert(2, twob)              # 2b -> position 2 (after intro)
ids = list(sldIdLst); fiveb = [x for x in ids if x is fiveb][0]
sldIdLst.remove(fiveb)
# conclusion is now last; insert 5b just before it
ids = list(sldIdLst); concl = ids[-1]; pos = list(sldIdLst).index(concl)
sldIdLst.insert(pos, fiveb)

out = os.path.join(ROOT, 'docs/Tanzania_SADC_LossAndDamage_Country_Statement.pptx')
p.save(out)
print('wrote', out, '|', len(p.slides._sldIdLst), 'slides')
